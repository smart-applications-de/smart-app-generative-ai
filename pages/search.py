# This pysqlite3 workaround is necessary for ChromaDB on Streamlit Sharing

import streamlit as st
import pandas as pd
import altair as alt
import datetime
import os
import io
import re
import markdown2
from docx import Document
from fpdf import FPDF
from bs4 import BeautifulSoup
import yfinance as yf
import plotly.express as px
import google.generativeai as gen
from google import genai
from langchain.agents import AgentType, Tool, initialize_agent
from langchain_community.utilities import GoogleSerperAPIWrapper
from langchain_google_genai import ChatGoogleGenerativeAI
from tools.utils import CrewStocknews

# --- Global Date/Year ---
date = pd.to_datetime('today').date()
yr = pd.to_datetime('today').year


# --- 1. Utility & File Conversion Functions ---

def convert_markdown_to_html(markdown_text):
    """Converts Markdown to HTML."""
    return markdown2.markdown(markdown_text, extras=["tables", "fenced-code-blocks", "strike"])


def convert_md_to_pdf(markdown_text):
    """Converts Markdown text to a PDF using FPDF, parsing the HTML structure."""
    try:
        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        for tag in soup.find_all(True):
            if tag.name == "h1":
                pdf.set_font("Arial", style="B", size=24)
                pdf.multi_cell(0, 10, txt=tag.get_text(), align="C")
                pdf.ln(5)
            elif tag.name == "h2":
                pdf.set_font("Arial", style="B", size=18)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(3)
            elif tag.name == "h3":
                pdf.set_font("Arial", style="B", size=14)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(2)
            elif tag.name == "p":
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(2)
            elif tag.name == "strong":
                pdf.set_font("Arial", style="B", size=12)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.set_font("Arial", size=12)
                pdf.ln(2)
            elif tag.name in ["ul", "ol"]:
                pdf.ln(2)
                for li in tag.find_all("li"):
                    pdf.set_font("Arial", size=12)
                    pdf.multi_cell(0, 10, txt=f"  •  {li.get_text()}")
                pdf.ln(2)
            elif tag.name == "hr":
                pdf.line(pdf.get_x(), pdf.get_y(), pdf.get_x() + 190, pdf.get_y())
                pdf.ln(5)

        return pdf.output(dest="S").encode("latin-1")
    except Exception as e:
        st.error(f"Error converting to PDF: {e}")
        return None


def convert_md_to_docx(markdown_text):
    """Converts Markdown text to a DOCX using python-docx, parsing HTML."""
    try:
        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")
        doc = Document()

        for tag in soup.find_all(True):
            if tag.name == "h1":
                doc.add_heading(tag.get_text(), level=1)
            elif tag.name == "h2":
                doc.add_heading(tag.get_text(), level=2)
            elif tag.name == "h3":
                doc.add_heading(tag.get_text(), level=3)
            elif tag.name == "p":
                doc.add_paragraph(tag.get_text())
            elif tag.name in ["ul", "ol"]:
                for li in tag.find_all("li"):
                    doc.add_paragraph(li.get_text(), style='List Bullet')
            elif tag.name == "hr":
                doc.add_paragraph("---")

        with io.BytesIO() as buffer:
            doc.save(buffer)
            return buffer.getvalue()
    except Exception as e:
        st.error(f"Error converting to DOCX: {e}")
        return None


def convert_md_to_pptx(markdown_text):
    """Converts Markdown to a basic PPTX."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders.get(1)

        title_text = soup.h1.text if soup.h1 else "AI Generated Report"
        title.text = title_text
        if subtitle and soup.h2:
            subtitle.text = soup.h2.text

        bullet_slide_layout = prs.slide_layouts[1]
        body_shape = None

        for tag in soup.find_all(['h3', 'p', 'ul', 'ol']):
            if tag.name == 'h3':
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                title_shape.text = tag.get_text()
            elif tag.name == 'p' and body_shape:
                p = body_shape.text_frame.add_paragraph()
                p.text = tag.get_text()
            elif tag.name in ['ul', 'ol'] and body_shape:
                for li in tag.find_all('li'):
                    p = body_shape.text_frame.add_paragraph()
                    p.text = li.get_text()
                    p.level = 1

        with io.BytesIO() as buffer:
            prs.save(buffer)
            return buffer.getvalue()
    except Exception as e:
        st.error(f"Error converting to PPTX: {e}")
        return None


def render_download_options(container, markdown_content, base_file_name, key_suffix):
    """Renders download buttons for various formats in a specified container."""
    if not markdown_content:
        container.warning("No content available to download.")
        return

    try:
        html_output = convert_markdown_to_html(markdown_content)
    except Exception as e:
        container.error(f"Error generating HTML for download: {e}")
        return

    download_format = container.radio(
        "Download as:",
        ("TXT", "PDF", "HTML", "MD", "DOCX", "PPTX"),
        key=f'download_{key_suffix}'
    )

    now = datetime.datetime.now()
    formatted_time = now.strftime("%Y-%m-%d_%H%M")
    file_name = f"{base_file_name}_{formatted_time}"

    try:
        if download_format == "TXT":
            container.download_button(
                label="Download TXT",
                data=markdown_content.encode("utf-8"),
                file_name=f"{file_name}.txt",
                mime="text/plain",
            )
        elif download_format == "MD":
            container.download_button(
                label="Download Markdown",
                data=markdown_content.encode("utf-8"),
                file_name=f"{file_name}.md",
                mime="text/markdown",
            )
        elif download_format == "HTML":
            container.download_button(
                label="Download HTML",
                data=html_output.encode("utf-8"),
                file_name=f"{file_name}.html",
                mime="text/html",
            )
        elif download_format == "PDF":
            pdf_bytes = convert_md_to_pdf(markdown_content)
            if pdf_bytes:
                container.download_button(
                    label="Download PDF",
                    data=pdf_bytes,
                    file_name=f"{file_name}.pdf",
                    mime="application/pdf",
                )
        elif download_format == "DOCX":
            docx_bytes = convert_md_to_docx(markdown_text)
            if docx_bytes:
                container.download_button(
                    label="Download DOCX",
                    data=docx_bytes,
                    file_name=f"{file_name}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
        elif download_format == "PPTX":
            pptx_bytes = convert_md_to_pptx(markdown_content)
            if pptx_bytes:
                container.download_button(
                    label="Download PPTX",
                    data=pptx_bytes,
                    file_name=f"{file_name}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
    except Exception as e:
        container.error(f"Error preparing download: {e}")


# --- 2. AI & Model Functions ---

@st.cache_data
def list_gemini_models():
    """Lists available Gemini models."""
    try:

        client = genai.Client(api_key=st.session_state['google_api_key'])
        gen.configure(api_key=st.session_state.germin_api_key)
        models = [
            m.name.split("/")[-1]
            for m in gen.list_models()
            if 'generateContent' in m.supported_generation_methods
        ]
        # Get all "pro" and "flash" models, etc.
        relevant_models = [m for m in models if "flash" in m or "pro" in m or "exp" in m or "2.5" in m]
        short_names = sorted(list(set(relevant_models)), reverse=True)
        full_paths = [f"gemini/{m}" for m in short_names]

        return short_names, full_paths
    except Exception as e:
        st.error(f"Error listing models: {e}")
        return [], []


@st.cache_resource
def SearchAgent(germin_key, SERPAPI_API_KEY, query, model="gemini-1.5-pro"):
    """Runs a general web search agent."""
    try:
        llm = ChatGoogleGenerativeAI(model=model, api_key=germin_key, temperature=0.3)
        search = GoogleSerperAPIWrapper(api_key=SERPAPI_API_KEY, num=10)
        tools = [
            Tool(
                name="Search",
                func=search.run,
                description="useful for when you need to ask with search.",
            )
        ]

        search_agent = initialize_agent(
            tools, llm, agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
            handle_parsing_errors=True, verbose=True
        )

        messages = [
            (
                "system",
                f"""You are a helpful assistant that with decades of experience on Google Search. 
                You take the user query and search on internet and returns the top results.
                You MUST Also Include source-information links for further reading.
                Formatted in markdown. The Current year is {yr} and current date is {date}.""",
            ),
            ("human", query),
        ]

        response = search_agent.run(messages)

        # Save to history
        st.session_state.searches.append({"role": "user", "content": query})
        st.session_state.searches.append({"role": "assistant", "content": response})
        return response
    except Exception as error:
        st.error(error)
        return f"An error occurred: {error}"


@st.cache_resource
def SearchNews(germin_key, SERPAPI_API_KEY, topic, model="gemini-1.5-pro"):
    """Runs a news-focused web search agent."""
    try:
        llm = ChatGoogleGenerativeAI(model=model, api_key=germin_key, temperature=0.3)
        news = GoogleSerperAPIWrapper(api_key=SERPAPI_API_KEY, num=20, type="news")
        tools_news = [
            Tool(
                name="news",
                func=news.run,
                description="useful for when you need to search for latest news.",
            )
        ]

        news_agent = initialize_agent(
            tools_news, llm, agent=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
            handle_parsing_errors=True, verbose=True
        )

        messages_news = [
            (
                "system",
                f"""You are a news Report Expert and you have decades of experience. 
                You take the user input on a given topic and search on internet for the latest news. 
                You MUST write an article with the most important and latest news. The current date is:{date} and current year:{yr}.
                You MUST Also Include source urls for further reading.
                Formatted in markdown.""",
            ),
            ("human", topic),
        ]
        response = news_agent.run(messages_news)
        return response
    except Exception as error:
        st.error(error)
        return f"An error occurred: {error}"


# --- 3. Main Application ---

# Check for API keys from session state (set in nav.py)
if 'google_api_key' not in st.session_state or not st.session_state['google_api_key']:
    st.error("🚫 Google Gemini API Key not found. Please set it in the main app.")
    st.stop()
if 'serper_api_key' not in st.session_state or not st.session_state['serper_api_key']:
    st.error("🚫 Serper API Key not found. Please set it in the main app.")
    st.stop()

# Configure API keys
GEMINI_API_KEY = st.session_state['google_api_key']
SERPER_API_KEY = st.session_state['serper_api_key']
os.environ['GOOGLE_API_KEY'] = GEMINI_API_KEY
os.environ['SERPER_API_KEY'] = SERPER_API_KEY
gen.configure(api_key=GEMINI_API_KEY)

# Get models
model_names, gemini_model_paths = list_gemini_models()
if not model_names:
    st.error("Could not fetch Gemini models. Check your API key.")
    st.stop()

# Initialize session state for chat history
if "searches" not in st.session_state:
    st.session_state["searches"] = [
        {"role": "assistant", "content": "Hi, I'm a chatbot who can search the web. How can I help you?"}
    ]

# --- UI Tabs ---
tab1, tab2, tab3 = st.tabs(["Agent Google Search", "Agent News Article", "Crew AI Stock News"])

with tab1:
    st.header("Agent-Powered Google Search")

    # Text input at the top
    query = st.text_input(
        "Ask Anything",
        key='query',
        placeholder="Who won the last F1 race?"
    )

    model1 = st.radio(
        "Choose a Model",
        model_names,  # Use short names for this agent
        key='model1',
        horizontal=True
    )

    if query:
        with st.spinner("Searching the web..."):
            response = SearchAgent(GEMINI_API_KEY, SERPER_API_KEY, query, model1)
            # The response is now automatically added to session state by the function

    st.write("---")
    st.subheader("Search History")
    # Display the full search history
    for msg in st.session_state.searches:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

with tab2:
    st.header("Agent-Powered News Search")

    news_topic = st.text_input(
        "Give a Topic",
        key='news_topic',
        placeholder="Latest AI advancements"
    )

    model2 = st.radio(
        "Choose a Model",
        model_names,  # Use short names for this agent
        key='model2',
        horizontal=True
    )

    if news_topic:
        with st.spinner(f"Searching for news on '{news_topic}'..."):
            response = SearchNews(GEMINI_API_KEY, SERPER_API_KEY, news_topic, model2)
            st.markdown(response)

with tab3:
    st.header("Crew AI Stock News Article")
    st.info("This agent will perform a deep dive to write a full article on a specific stock.")

    container3 = st.container(border=True)

    topic = container3.text_input(
        "News topic (Stock Ticker)",
        key='topic',
        placeholder="e.g., NVDA or AAPL"
    )

    model3 = container3.radio(
        "Choose a Model",
        gemini_model_paths,  # Crew AI needs the full path
        key='model3'
    )

    if container3.button("Generate Stock News Article"):
        if topic and model3:
            news_article = None
            with st.spinner(f"AI Crew is writing an article on '{topic}'... This may take several minutes."):
                try:
                    # Note: CrewStocknews is assumed to not take company, sector, summary
                    results = CrewStocknews(
                        germin_api=GEMINI_API_KEY,
                        serp_api=SERPER_API_KEY,
                        topic=topic,
                        year=yr,
                        date=date,
                        model=model3
                    )
                    st.markdown(results)

                    # Try to read the generated file
                    file_name = f'{topic}_editor_task.md'  # Assuming this is the output file
                    if os.path.exists(file_name):
                        with open(file_name, 'r', encoding='utf-8') as f:
                            news_article = f.read()
                        st.session_state.last_crew_article = news_article
                        st.session_state.last_crew_topic = topic
                    else:
                        st.warning(f"Crew finished, but output file '{file_name}' not found. Showing main output only.")
                        st.session_state.last_crew_article = results  # Fallback
                        st.session_state.last_crew_topic = topic

                except Exception as e:
                    st.error(f"An error occurred with the AI Crew: {e}")
        else:
            container3.warning("Please provide a stock ticker and select a model.")

    # Download section - shows if an article is in session state
    if "last_crew_article" in st.session_state:
        st.subheader(f"Download Article for {st.session_state.last_crew_topic}")
        render_download_options(
            st,
            st.session_state.last_crew_article,
            f"{st.session_state.last_crew_topic}_ai_news_article",
            "crew_news_download"
        )
