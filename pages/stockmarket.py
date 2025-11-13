import pandas as pd
import altair as alt
import streamlit as st
import datetime
import os
import io
import re
import markdown2
from docx import Document
from fpdf import FPDF
from bs4 import BeautifulSoup
import google.generativeai as gen
from google import genai
import yfinance as yf
import plotly.express as px

# --- 1. Utility & File Conversion Functions ---

@st.cache_data
def load_csv_data(file_path):
    """A cached function to load CSV data."""
    if not os.path.exists(file_path):
        st.error(f"Error: Data file not found at {file_path}")
        return pd.DataFrame()
    try:
        return pd.read_csv(file_path)
    except Exception as e:
        st.error(f"Error loading {file_path}: {e}")
        return pd.DataFrame()


def convert_to_billions(value):
    """
    Converts a string value (e.g., '4.93T', '150B', '500M') to a float in billions.
    """
    if pd.isna(value):
        return pd.NA

    value_str = str(value).strip().upper().replace(',', '')  # Clean the string

    if value_str.endswith('T'):
        # Convert Trillions to Billions
        return pd.to_numeric(value_str[:-1], errors='coerce') * 1000
    elif value_str.endswith('B'):
        # Convert Billions to Billions
        return pd.to_numeric(value_str[:-1], errors='coerce')
    elif value_str.endswith('M'):
        # Convert Millions to Billions
        return pd.to_numeric(value_str[:-1], errors='coerce') / 1000
    else:
        # Coerce errors for any non-conforming strings
        return pd.to_numeric(value_str, errors='coerce')


def convert_markdown_to_html(markdown_text):
    """Converts Markdown to HTML."""
    return markdown2.markdown(markdown_text)


def convert_md_to_pdf(markdown_text):
    """Converts Markdown text to a PDF using FPDF, parsing the HTML structure."""
    try:
        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        for tag in soup.find_all(True):
            if tag.name == "h1":
                pdf.set_font("Arial", style="B", size=24)
                pdf.multi_cell(0, 10, txt=tag.get_text(), align="C")
                pdf.ln(5)
            elif tag.name == "h2":
                pdf.set_font("Arial", style="B", size=18)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(3)
            elif tag.name == "h3":
                pdf.set_font("Arial", style="B", size=14)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(2)
            elif tag.name == "p":
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.ln(2)
            elif tag.name == "strong":
                pdf.set_font("Arial", style="B", size=12)
                pdf.multi_cell(0, 10, txt=tag.get_text())
                pdf.set_font("Arial", size=12)
                pdf.ln(2)
            elif tag.name in ["ul", "ol"]:
                pdf.ln(2)
                for li in tag.find_all("li"):
                    pdf.set_font("Arial", size=12)
                    pdf.multi_cell(0, 10, txt=f"  •  {li.get_text()}")
                pdf.ln(2)
            elif tag.name == "hr":
                pdf.line(pdf.get_x(), pdf.get_y(), pdf.get_x() + 190, pdf.get_y())
                pdf.ln(5)

        return pdf.output(dest="S").encode("latin-1")
    except Exception as e:
        st.error(f"Error converting to PDF: {e}")
        return None


def convert_md_to_docx(markdown_text):
    """Converts Markdown text to a DOCX using python-docx, parsing HTML."""
    try:
        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")
        doc = Document()

        for tag in soup.find_all(True):
            if tag.name == "h1":
                doc.add_heading(tag.get_text(), level=1)
            elif tag.name == "h2":
                doc.add_heading(tag.get_text(), level=2)
            elif tag.name == "h3":
                doc.add_heading(tag.get_text(), level=3)
            elif tag.name == "p":
                doc.add_paragraph(tag.get_text())
            elif tag.name in ["ul", "ol"]:
                for li in tag.find_all("li"):
                    doc.add_paragraph(li.get_text(), style='List Bullet')
            elif tag.name == "hr":
                doc.add_paragraph("---")

        with io.BytesIO() as buffer:
            doc.save(buffer)
            return buffer.getvalue()
    except Exception as e:
        st.error(f"Error converting to DOCX: {e}")
        return None


def convert_md_to_pptx(markdown_text):
    """Converts Markdown to a basic PPTX."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        html = convert_markdown_to_html(markdown_text)
        soup = BeautifulSoup(html, "html.parser")

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders.get(1)

        title_text = soup.h1.text if soup.h1 else "AI Generated Report"
        title.text = title_text
        if subtitle and soup.h2:
            subtitle.text = soup.h2.text

        bullet_slide_layout = prs.slide_layouts[1]
        body_shape = None

        for tag in soup.find_all(['h3', 'p', 'ul', 'ol']):
            if tag.name == 'h3':
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                title_shape.text = tag.get_text()
            elif tag.name == 'p' and body_shape:
                p = body_shape.text_frame.add_paragraph()
                p.text = tag.get_text()
            elif tag.name in ['ul', 'ol'] and body_shape:
                for li in tag.find_all('li'):
                    p = body_shape.text_frame.add_paragraph()
                    p.text = li.get_text()
                    p.level = 1

        with io.BytesIO() as buffer:
            prs.save(buffer)
            return buffer.getvalue()
    except Exception as e:
        st.error(f"Error converting to PPTX: {e}")
        return None


def render_download_options(container, markdown_content, base_file_name, key_suffix):
    """Renders download buttons for various formats in a specified container."""
    if not markdown_content:
        container.warning("No content available to download.")
        return

    try:
        html_output = convert_markdown_to_html(markdown_content)
    except Exception as e:
        container.error(f"Error generating HTML for download: {e}")
        return

    download_format = container.radio(
        "Download as:",
        ("TXT", "PDF", "HTML", "MD", "DOCX", "PPTX"),
        key=f'download_{key_suffix}'
    )

    now = datetime.datetime.now()
    formatted_time = now.strftime("%Y-%m-%d_%H%M")
    file_name = f"{base_file_name}_{formatted_time}"

    try:
        if download_format == "TXT":
            container.download_button(
                label="Download TXT",
                data=html_output,
                file_name=f"{file_name}.txt",
                mime="text/plain",
            )
        elif download_format == "MD":
            container.download_button(
                label="Download Markdown",
                data=html_output,
                file_name=f"{file_name}.md",
                mime="text/markdown",
            )

    except Exception as e:
        container.error(f"Error preparing download: {e}")


# --- 2. API Key Management (Sidebar) ---

def setup_api_keys():
    """Manages API key input in the sidebar and returns True if keys are set."""
    st.sidebar.header("⚙️ API Configuration")

    if "google_api_key" not in st.session_state:
        st.session_state.google_api_key = ""
    if "serper_api_key" not in st.session_state:
        st.session_state.serper_api_key = ""

    google_key = st.sidebar.text_input(
        "Google Gemini API Key",
        key="google_api_key_input",
        type="password",
        value=st.session_state.google_api_key
    )
    st.sidebar.markdown("[Get Google API Key](https://ai.google.dev/)")

    serper_key = st.sidebar.text_input(
        "Serper API Key",
        key="serper_api_key_input",
        type="password",
        value=st.session_state.serper_api_key
    )
    st.sidebar.markdown("[Get Serper API Key](https://serper.dev/)")

    if google_key and serper_key:
        st.session_state.google_api_key = google_key
        st.session_state.serper_api_key = serper_key
        return True

    st.info("Please enter your Google Gemini and Serper API keys in the sidebar to start the analysis.")
    return False


# --- Main Application ---
if setup_api_keys():

    # --- 3. Conditional Imports & AI Setup ---
    GOOGLE_API_KEY = st.session_state.google_api_key
    SERPER_API_KEY = st.session_state.serper_api_key

    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    os.environ['SERPER_API_KEY'] = SERPER_API_KEY
    os.environ["OPENAI_API_KEY"] = GOOGLE_API_KEY
    try:
        gen.configure(api_key=GOOGLE_API_KEY)
        client = genai.Client(api_key=GOOGLE_API_KEY)
    except Exception as e:
        st.error(f"Failed to configure Google Gemini AI: {e}")
        st.stop()

    try:
        from history.data import DailyQuote, AllStocksPerformnace
        from history.webscraping import SummaryLoad
        from Forecasts.models import VECM_Forecast_Lower_Upper
        from StockmarketAnalysis.crew import StockMarketCrew
        from News.news_and_marketing import CrewStocknews
    except ImportError as e:
        st.error(f"Failed to import necessary modules: {e}")
        st.error("Please ensure all required packages (history, Forecasts, etc.) are in the correct path.")
        st.stop()





    @st.cache_data
    def list_gemini_models(_client):
        """Lists available Gemini models."""
        try:
            models = [
                m.name.split("/")[-1]
                for m in gen.list_models()
                if 'generateContent' in m.supported_generation_methods
            ]
            vision_models = [f"gemini/{m}" for m in models if
                             "2.0" in m.lower() or "-exp" in m or "2.5" in m or "1.5" in m]
            return vision_models
        except Exception as e:
            st.error(f"Error listing models: {e}")
            return []


    gemini_models = list_gemini_models(client)


    # --- 4. Cached Data Loading Functions ---

    @st.cache_data
    def get_all_tickers():
        """Consolidates tickers from multiple sources."""
        try:
            df_mega = load_csv_data("./history/data/mega_capCompanies.csv")
            df_jones = load_csv_data('./history/data/dow_jones.csv')
            df_dax = load_csv_data("./history/data/dax_companies.csv")
            df_symbol_txt = load_csv_data("./history/data/symbol.txt")

            df_symbol = pd.DataFrame()
            if not df_mega.empty:
                df_symbol = pd.concat([df_symbol, df_mega[['Symbol']]], ignore_index=True)
            if not df_jones.empty:
                df_symbol = pd.concat([df_symbol, df_jones[['Symbol']]], ignore_index=True)
            if not df_dax.empty:
                sy = df_dax[['Ticker-en']].rename(columns={'Ticker-en': 'Symbol'})
                df_symbol = pd.concat([df_symbol, sy], ignore_index=True)
            if not df_symbol_txt.empty:
                per = df_symbol_txt.rename(columns={df_symbol_txt.columns[0]: 'Symbol'})
                df_symbol = pd.concat([df_symbol, per], ignore_index=True)

            df_symbol = df_symbol.drop_duplicates(ignore_index=True)
            return df_symbol
        except Exception as e:
            st.error(f"Error compiling tickers: {e}")
            return pd.DataFrame()


    @st.cache_data
    def get_stock_performance():
        """Loads and processes performance data for all tickers."""
        try:
            df_stock = get_all_tickers().rename(columns={'Symbol': 'symbol'})
            if df_stock.empty:
                return pd.DataFrame()

            # Assuming AllStocksPerformnace returns a 'sector' column
            df = AllStocksPerformnace(df_stock).reset_index(drop=True)
            df = df.sort_values(by='1yr_performance', ascending=False).reset_index(drop=True)
            return df
        except Exception as e:
            st.error(f"Error getting stock performance: {e}")
            return pd.DataFrame()


    @st.cache_data
    def get_historical_data(symbol):
        try:
            df, df_monthly, df_data, df_f = DailyQuote(symbol)
            return df, df_monthly, df_data, df_f
        except Exception as e:
            st.error(f"Error getting historical data for {symbol}: {e}")
            return pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame()


    @st.cache_resource
    def get_forecast_data(_df):
        """Generates VECM forecast."""
        if _df.empty:
            return pd.DataFrame()
        try:
            df_sorted = (_df.sort_values('date', ascending=True)).reset_index(drop=True)
            max_date = (df_sorted['date'].dt.date.max() + pd.DateOffset(days=1))
            pred = VECM_Forecast_Lower_Upper(
                df_sorted, steps=20, seasons=252,
                response="close", start_date=max_date, freq='B'
            )
            return pred
        except Exception as e:
            st.error(f"Error generating forecast: {e}")
            return pd.DataFrame()


    # --- 5. AI Crew Functions ---

    @st.cache_resource
    def get_crew_ai_news(topic, model):
        """Runs the AI news crew."""
        try:
            date = pd.to_datetime('today').date()
            yr = pd.to_datetime('today').year
            p = yf.Ticker(str(topic))
            info = p.info
            company = info.get('longName', 'N/A')
            sector = info.get('sector', 'N/A')
            summary = info.get('longBusinessSummary', 'N/A')
            results = CrewStocknews(GOOGLE_API_KEY, SERPER_API_KEY, topic, yr, date, company, sector, summary,
                                    str(model))
            return results
        except Exception as e:
            st.error(f"Error running AI News Crew: {e}")
            return None


    @st.cache_resource
    def get_crew_ai_recommendation(option, model):
        """Runs the AI stock recommendation crew."""
        try:
            p = yf.Ticker(str(option))
            info = p.info
            company = info.get('longName', 'N/A')
            sector = info.get('sector', 'N/A')
            summary = info.get('longBusinessSummary', 'N/A')
            industry = info.get('industry', 'N/A')
            website = info.get('website', 'N/A')

            history = f'monthly_{option}_data.csv'
            financial_data = f'_{str(option).lower()}_qfinancial_ratios.csv'
            company_stock_info = f'company_stock_info_{option}_data'

            results = StockMarketCrew(
                GOOGLE_API_KEY, SERPER_API_KEY, option, history, financial_data,
                company, sector, summary, industry, website, company_stock_info, str(model)
            ).run()
            return results
        except Exception as e:
            st.error(f"Error running AI Recommendation Crew: {e}")
            return None


    # --- 6. Streamlit UI (Tabs) ---

    st.title("📈 Stockmarket Analysis and Financial Advisor with AI")

    all_tickers_df = get_all_tickers()
    if all_tickers_df.empty:
        st.error("Could not load any stock tickers. Please check data files.")
        st.stop()

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "Overview", "Stock Performance", "Individual Stock", "AI News Article", "AI Stock Recommendation"
    ])

    with tab1:
        st.header("Industry Performance")
        try:
            df1 = load_csv_data("./history/data/industry.csv")
            if not df1.empty:
                df1 = df1.replace({'%': ''}, regex=True)
                df1['Market Cap'] = df1['Market Cap'].apply(convert_to_billions)
                df1['Profit Margin'] = pd.to_numeric(df1['Profit Margin'], errors='coerce')
                df1['1D Change'] = pd.to_numeric(df1['1D Change'], errors='coerce')
                df1['1Y Change'] = pd.to_numeric(df1['1Y Change'], errors='coerce')

                col1, col2, col3 = st.columns(3)
                col1.metric("Max 1D Change", f"{df1['1D Change'].max():.2f}%")
                col2.metric("Max Market Cap (Billions)", f"${df1['Market Cap'].max():.2f}B")
                col3.metric("No. Industries", df1['Industry Name'].count())

                st.dataframe(df1, use_container_width=True)

                # --- VIZ UPGRADE: Replaced Pie with Sorted Horizontal Bar ---
                chart_data_ind = df1.dropna(subset=['1D Change']).nlargest(20, '1D Change')
                bar_ind = alt.Chart(chart_data_ind, title="Top 20 Industries by 1-Day Performance").mark_bar().encode(
                    x=alt.X('1D Change:Q', title='1-Day Change (%)'),
                    y=alt.Y('Industry Name:N', sort='-x'),  # Sort descending
                    tooltip=['Industry Name', '1D Change']
                ).interactive()
                st.altair_chart(bar_ind, theme="streamlit", use_container_width=True)
        except Exception as e:
            st.error(f"Error in Industry Performance: {e}")

        st.header("Sector Performance")
        try:
            df2 = load_csv_data("./history/data/sector.csv")
            if not df2.empty:
                df2 = df2.replace({',': '', '%': ''}, regex=True)
                df2['Market Cap'] = df2['Market Cap'].apply(convert_to_billions)
                df2['1D Change'] = pd.to_numeric(df2['1D Change'], errors='coerce')
                df2['1Y Change'] = pd.to_numeric(df2['1Y Change'], errors='coerce')

                st.dataframe(df2, use_container_width=True)

                # --- VIZ UPGRADE: Replaced Pie with Sorted Horizontal Bar ---
                chart_data_sec = df2.dropna(subset=['1Y Change'])
                bar_sec = alt.Chart(chart_data_sec, title="Sector 1-Year Performance").mark_bar().encode(
                    x=alt.X('1Y Change:Q', title='1-Year Change (%)'),
                    y=alt.Y('Sector Name:N', sort='-x'),  # Sort descending
                    color='Sector Name:N',
                    tooltip=['Sector Name', '1Y Change']
                ).interactive()
                st.altair_chart(bar_sec, theme="streamlit", use_container_width=True)
        except Exception as e:
            st.error(f"Error in Sector Performance: {e}")

    with tab2:
        st.header("Stock Performance Dashboard")
        try:
            df_per = get_stock_performance()
            if not df_per.empty:

                # --- FILTER UPGRADE: Added Sector and Performance Filters ---
                filt_col1, filt_col2 = st.columns(2)

                # Sector Filter (check if 'sector' column exists)
                if 'sector' in df_per.columns:
                    sectors = sorted(df_per['sector'].dropna().unique())
                    selected_sectors = filt_col1.multiselect('Filter by Sector', sectors, default=sectors)
                else:
                    selected_sectors = []
                    filt_col1.info("Sector data not available for filtering.")

                # Performance Filter
                min_val = float(df_per['1yr_performance'].min())
                max_val = float(df_per['1yr_performance'].max())
                default_val = float(df_per['1yr_performance'].quantile(0.1))  # Start by filtering bottom 10%

                min_perf = filt_col2.slider(
                    'Filter by 1-Year Performance >',
                    min_value=min_val,
                    max_value=max_val,
                    value=default_val,
                    format="%.2f%%"
                )

                # Apply filters
                if 'sector' in df_per.columns and selected_sectors:
                    filtered_df_per = df_per[
                        (df_per['sector'].isin(selected_sectors)) &
                        (df_per['1yr_performance'] > min_perf)
                        ]
                else:
                    filtered_df_per = df_per[df_per['1yr_performance'] > min_perf]

                st.dataframe(filtered_df_per, column_config={
                    "website": st.column_config.LinkColumn(),
                    "Date": st.column_config.DateColumn(),
                }, use_container_width=True)
            else:
                st.warning("Could not load stock performance data.")
        except Exception as perf:
            st.error(perf)

        st.subheader("ETFs and Indices Performance")
        try:
            df_etf = load_csv_data("./history/data/df_etf.csv")
            if not df_etf.empty:
                df_etf = df_etf.sort_values(by='1yr_performance', ascending=False).reset_index(drop=True)
                st.dataframe(df_etf, use_container_width=True)

                # --- VIZ UPGRADE: Sorted Bar Chart ---
                bareft = alt.Chart(df_etf, title="EFT'S and Indices Performance").mark_bar().encode(
                    x=alt.X('Ticker:N', sort='-y'),  # Sort by performance
                    y='1yr_performance:Q',
                    color='Ticker:N',
                    tooltip=['Ticker', '1yr_performance']
                ).interactive()
                st.altair_chart(bareft, theme="streamlit", use_container_width=True)
        except Exception as etf:
            st.error(etf)

        st.header("Top S&P 500 Companies")
        try:
            df5 = load_csv_data("./history/data/s_p500.csv")
            if not df5.empty:
                # --- FILTER UPGRADE: Added GICS Sector Filter ---
                gics_sectors = sorted(df5['GICS Sector'].dropna().unique())
                selected_gics = st.multiselect('Filter by GICS Sector', gics_sectors, default=gics_sectors)

                filtered_df5 = df5[df5['GICS Sector'].isin(selected_gics)]
                st.dataframe(filtered_df5, use_container_width=True)

                # Chart now uses the filtered data
                df_sectory = filtered_df5.groupby('GICS Sector')['Symbol'].count().reset_index().rename(
                    columns={'Symbol': 'Counts'})

                # --- VIZ UPGRADE: Sorted Bar Chart ---
                bar4 = alt.Chart(df_sectory, title="Top S&P 500 Companies by Sector (Filtered)").mark_bar().encode(
                    x=alt.X('GICS Sector', sort='-y'),  # Sort by count
                    y='Counts:Q',
                    color='GICS Sector:N',
                    tooltip=['GICS Sector', 'Counts']
                ).interactive()
                st.altair_chart(bar4, theme="streamlit", use_container_width=True)
        except Exception as e:
            st.error(f"Error displaying S&P 500 data: {e}")

    with tab3:
        st.header("Individual Stock Analysis")
        option = st.selectbox(
            "Which Stock Would you like to analyze?",
            all_tickers_df['Symbol'],
            index=all_tickers_df['Symbol'].tolist().index('AAPL') if 'AAPL' in all_tickers_df['Symbol'].tolist() else 0,
            key='stock_select'
        )

        if option:
            try:
                df, df_monthly, df_data, df_f = get_historical_data(option)

                st.subheader("Daily Historical Data")
                if not df.empty:
                    st.dataframe(df, use_container_width=True, column_config={"date": st.column_config.DateColumn()})
                    fig = px.line(
                        df, x="date", y=['close', 'sma_5_day', 'sma_20_day', 'sma_50_day', 'sma_200_day'],
                        hover_data={"date": "|%B %d, %Y"},
                        title=f'{option} Stock Price Evolution'
                    )
                    fig.update_xaxes(dtick="M1", tickformat="%Y %B", ticklabelmode="period")
                    st.plotly_chart(fig, theme="streamlit", use_container_width=True)

                st.subheader("Fundamental Analysis - Company Information")
                if not df_data.empty:
                    st.dataframe(df_data, use_container_width=True,
                                 column_config={"website": st.column_config.LinkColumn()})

                st.subheader("Fundamental Analysis - Company Financial Ratios")
                if not df_f.empty:
                    st.dataframe(df_f, use_container_width=True)

                st.header("Forecast with VECM Model")
                pred = get_forecast_data(df)
                if not pred.empty:
                    st.dataframe(pred, use_container_width=True,
                                 column_config={"prediction_date": st.column_config.DateColumn()})
                    fig1 = px.line(
                        pred, x="prediction_date", y=['close', 'close_lower', 'close_upper'],
                        hover_data={"prediction_date": "|%B %d, %Y"},
                        title=f'{option} Stock Price Forecast (Next 20 Days)'
                    )
                    fig1.update_xaxes(tickformat="%d-%m-%Y")
                    st.plotly_chart(fig1, theme="streamlit", use_container_width=True)
                else:
                    st.warning("Could not generate forecast for this stock.")

            except Exception as errordf:
                st.error(f"An error occurred analyzing {option}: {errordf}")

    with tab4:
        st.header("📰 AI-Generated Stock News Article")
        container = st.container(border=True)

        col1, col2 = container.columns(2)
        model1 = col1.radio(
            "Choose a Model (News)",
            gemini_models,
            key='model1',
            index=0 if gemini_models else -1
        )
        option1 = col2.selectbox(
            "Which Stock do you want news for?",
            all_tickers_df['Symbol'],
            key='option1',
            index=all_tickers_df['Symbol'].tolist().index('MSFT') if 'MSFT' in all_tickers_df['Symbol'].tolist() else 0
        )

        if container.button("Generate News Article", key="gen_news"):
            if model1 and option1:
                with st.spinner(f"Generating news for {option1}... This may take a moment."):
                    ai_news_report = get_crew_ai_news(option1, model1)
                    if ai_news_report:
                        st.session_state.ai_news_report = ai_news_report
                        st.session_state.current_news_stock = option1



                    else:
                        st.error("Failed to generate news report.")
            else:
                st.warning("Please select a model and a stock.")

        if "ai_news_report" in st.session_state:
            st.subheader(f"AI-Generated Article for {st.session_state.current_news_stock}")
            st.markdown(st.session_state.ai_news_report)

            st.subheader("Download News Article")
            file_news = f'{option1}_ai_news_write_task.md'
            with open(file_news, 'r') as f:
                news_expert = (f.read())
                f.close()
            render_download_options(
                st,
                news_expert,
                f"{st.session_state.current_rec_stock}_ai_report",
                "news"
            )


    with tab5:
        st.header("🤖 AI-Agent Stock Recommendation")
        container1 = st.container(border=True)

        col3, col4 = container1.columns(2)
        model2 = col3.radio(
            "Choose a Model (Recommendation)",
            gemini_models,
            key='model2',
            index=0 if gemini_models else -1
        )
        option2 = col4.selectbox(
            "Which Stock do you want a recommendation for?",
            all_tickers_df['Symbol'],
            key='option2',
            index=all_tickers_df['Symbol'].tolist().index('NVDA') if 'NVDA' in all_tickers_df['Symbol'].tolist() else 0
        )

        if container1.button("Generate Recommendation", key="gen_rec"):
            if model2 and option2:
                with st.spinner(f"Generating recommendation for {option2}... This may take several minutes."):
                    ai_recommendation = get_crew_ai_recommendation(option2, model2)
                    if ai_recommendation:
                        st.session_state.ai_recommendation = ai_recommendation
                        st.session_state.current_rec_stock = option2
                    else:
                        st.error("Failed to generate recommendation.")
            else:
                st.warning("Please select a model and a stock.")

        if "ai_recommendation" in st.session_state:
            st.subheader(f"AI-Generated Report for {st.session_state.current_rec_stock}")
            st.markdown(st.session_state.ai_recommendation)

            # Read and display other generated markdown files
            try:
                file_st = f'./Crew_AI/Reports/{option2}_filereader_task.md'
                with open(file_st, 'r') as f:
                    st.subheader("Reader and Fundamental Analyst")
                    st.markdown(f.read())
            except FileNotFoundError:
                pass  # It's okay if this file doesn't exist
            except Exception as e:
                st.warning(f"Could not read filereader task: {e}")

            try:
                file_f = f'./Crew_AI/Reports/{option2}_financial_analysis_task.md'
                with open(file_f, 'r') as f:
                    st.subheader("Finance Analyst")
                    st.markdown(f.read())
            except FileNotFoundError:
                pass
            except Exception as e:
                st.warning(f"Could not read financial analysis task: {e}")

            try:
                file_r = f'./Crew_AI/Reports/{option2}_research_task.md'
                with open(file_r, 'r') as f:
                    st.subheader("News and Market Sentiment Analyst")
                    advisor_expert = (f.read())
                    st.markdown(advisor_expert)
            except FileNotFoundError:
                pass
            except Exception as e:
                st.warning(f"Could not read research task: {e}")

            st.subheader("Download Full Investment Report")
            # Use the main recommendation output for download
            render_download_options(
                st,
                advisor_expert,
                f"{st.session_state.current_rec_stock}_ai_report",
                "recommendation"
            )

else:
    st.warning("Please enter your API keys in the sidebar to load the application.")