import os
import io
import re
import datetime
import streamlit as st
import markdown2
from bs4 import BeautifulSoup
from docx import Document
from fpdf import FPDF
from pytubefix import YouTube
from langchain_community.document_loaders import PyPDFLoader
#import google.generativeai as genai
from google import genai
from google.genai import types 
import google.generativeai as gen

# # --- Page Configuration ---
# st.set_page_config(
#     page_title="Chat with Gemini",
#     page_icon="💬",
#     layout="wide"
# )

# --- Session State Initialization ---
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "How can I help you today?"}]
if "gemini_client" not in st.session_state:
    st.session_state.gemini_client = None
if "context" not in st.session_state:
    st.session_state.context = None

# --- 1. API Key and Model Helpers ---

@st.cache_data
def list_gemini_models():
    """Lists available Gemini models that support content generation."""
    try:
        models = [
            m.name.split("/")[-1] 
            for m in gen.list_models()
            if 'generateContent' in m.supported_generation_methods
        ]
        # Prioritize flash and pro models
        vision_models = sorted([m for m in models if "flash" in m or "pro" in m], reverse=True)
        return vision_models
    except Exception as e:
        st.error(f"Error listing models: {e}")
        return []

def get_gemini_api_key():
    """Gets the Gemini API key from the user and initializes the client."""
    st.sidebar.header("⚙️ Configuration")
    if 'germin_api_key' in st.session_state and st.session_state.germin_api_key:
        try:
            #genai.configure(api_key=st.session_state.germin_api_key)
            st.session_state.gemini_client = genai.Client(api_key=st.session_state.germin_api_key)
            gen.configure(api_key=st.session_state.germin_api_key)
            st.session_state.model_client = genai.Client(api_key=st.session_state.germin_api_key)

            return st.session_state.germin_api_key
        except Exception as e:
            st.sidebar.error(f"Invalid API Key: {e}")
            del st.session_state.germin_api_key # Clear invalid key
            st.session_state.gemini_client = None
            
    google_api_key = st.sidebar.text_input(
        "Google Gemini API Key", key="gemini_api_key_input", type="password"
    )
    st.sidebar.markdown("[Get your Google API Key](https://ai.google.dev/)")

    if google_api_key:
        st.session_state['germin_api_key'] = google_api_key
        try:
            #genai.configure(api_key=google_api_key)
            st.session_state.gemini_client = genai.Client(api_key=google_api_key)
            st.rerun()
        except Exception as e:
            st.sidebar.error(f"Invalid API Key: {e}")
            del st.session_state.germin_api_key
            st.session_state.gemini_client = None
    
    st.info("Please enter your Google Gemini API Key in the sidebar to begin.")
    return None

# --- 2. Chat and Response Helpers ---

def handle_chat_response(prompt, response):
    """Adds user prompt and AI response to session state and displays them."""
    st.session_state.messages.append({"role": "user", "content": prompt})
    st.session_state.messages.append({"role": "assistant", "content": response.text})

def generate_gemini_content(client, model, contents):
    """Generates content using the Gemini API."""
    try:
        response = client.models.generate_content(model=model, contents=contents)
        return response
    except Exception as e:
        st.error(f"Error generating content: {e}")
        return None

def format_chat_history(messages):
    """Formats the chat history for download."""
    md_text = ""
    for msg in messages:
        if msg["role"] == "user":
            md_text += f"**You:**\n{msg['content']}\n\n---\n\n"
        else:
            md_text += f"**Gemini:**\n{msg['content']}\n\n---\n\n"
    return md_text

# --- 3. File Conversion Helpers ---

def convert_markdown_to_html(markdown_text):
    """Converts Markdown to HTML."""
    return markdown2.markdown(markdown_text, extras=["tables", "fenced-code-blocks", "strike"])

def convert_md_to_pdf(markdown_text):
    """Converts Markdown text to a PDF using FPDF, parsing the HTML structure."""
    html = convert_markdown_to_html(markdown_text)
    soup = BeautifulSoup(html, "html.parser")
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    for tag in soup.find_all(True):
        if tag.name == "h1":
            pdf.set_font("Arial", style="B", size=24)
            pdf.multi_cell(0, 10, txt=tag.get_text(), align="C")
            pdf.ln(5)
        elif tag.name == "h2":
            pdf.set_font("Arial", style="B", size=18)
            pdf.multi_cell(0, 10, txt=tag.get_text())
            pdf.ln(3)
        elif tag.name == "h3":
            pdf.set_font("Arial", style="B", size=14)
            pdf.multi_cell(0, 10, txt=tag.get_text())
            pdf.ln(2)
        elif tag.name == "p":
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, txt=tag.get_text())
            pdf.ln(2)
        elif tag.name == "strong":
            pdf.set_font("Arial", style="B", size=12)
            pdf.multi_cell(0, 10, txt=tag.get_text())
            pdf.set_font("Arial", size=12)
            pdf.ln(2)
        elif tag.name in ["ul", "ol"]:
            pdf.ln(2)
            for li in tag.find_all("li"):
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, txt=f"  •  {li.get_text()}")
            pdf.ln(2)
        elif tag.name == "hr":
            pdf.line(pdf.get_x(), pdf.get_y(), pdf.get_x() + 190, pdf.get_y())
            pdf.ln(5)

    return pdf.output(dest="S").encode("latin-1")

def convert_md_to_docx(markdown_text):
    """Converts Markdown text to a DOCX using python-docx, parsing HTML."""
    html = convert_markdown_to_html(markdown_text)
    soup = BeautifulSoup(html, "html.parser")
    doc = Document()

    for tag in soup.find_all(True):
        if tag.name == "h1":
            doc.add_heading(tag.get_text(), level=1)
        elif tag.name == "h2":
            doc.add_heading(tag.get_text(), level=2)
        elif tag.name == "h3":
            doc.add_heading(tag.get_text(), level=3)
        elif tag.name == "p":
            doc.add_paragraph(tag.get_text())
        elif tag.name in ["ul", "ol"]:
            for li in tag.find_all("li"):
                doc.add_paragraph(li.get_text(), style='List Bullet')
        elif tag.name == "hr":
            doc.add_paragraph("---")
            
    # Save to in-memory buffer for download
    with io.BytesIO() as buffer:
        doc.save(buffer)
        return buffer.getvalue()

def convert_md_to_pptx(markdown_text):
    """Converts Markdown to a basic PPTX."""
    from pptx import Presentation
    from pptx.util import Inches
    
    html = convert_markdown_to_html(markdown_text)
    soup = BeautifulSoup(html, "html.parser")
    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders.get(1)

    title_text = soup.h1.text if soup.h1 else "Chat History"
    title.text = title_text
    if subtitle and soup.h2:
        subtitle.text = soup.h2.text
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Add content slides (basic paragraphs and lists)
    current_slide = None
    body_shape = None

    for tag in soup.find_all(['h3', 'p', 'ul', 'ol']):
        if tag.name == 'h3':
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = current_slide.shapes.title
            body_shape = current_slide.placeholders[1]
            title_shape.text = tag.get_text()
        elif tag.name == 'p' and body_shape:
            p = body_shape.text_frame.add_paragraph()
            p.text = tag.get_text()
        elif tag.name in ['ul', 'ol'] and body_shape:
            for li in tag.find_all('li'):
                p = body_shape.text_frame.add_paragraph()
                p.text = li.get_text()
                p.level = 1

    with io.BytesIO() as buffer:
        prs.save(buffer)
        return buffer.getvalue()

# --- 4. Media Processing Helpers ---

def get_youtube_stream(video_url, stream_type="audio"):
    """Gets the audio or video stream from a YouTube URL."""
    try:
        yt = YouTube(video_url)
        if stream_type == "audio":
            return yt.streams.filter(only_audio=True).first()
        else:
            return yt.streams.get_highest_resolution()
    except Exception as e:
        st.error(f"Error processing YouTube URL: {e}")
        return None

def process_pdf(file_bytes):
    """Loads and splits a PDF, returning its text content."""
    try:
        # Save to a temporary file
        if not os.path.isdir("temp_dir"):
            os.mkdir("temp_dir")
        path = os.path.join("temp_dir", "temp.pdf")
        
        with open(path, "wb") as f:
            f.write(file_bytes)
            
        loader = PyPDFLoader(path)
        pages = loader.load_and_split()
        article = " ".join([doc.page_content for doc in pages])
        os.remove(path) # Clean up
        return article
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return None

# --- 5. Main Application UI ---

# Get API Key and initialize client
api_key = get_gemini_api_key()

if not api_key or not st.session_state.gemini_client:
    st.stop()

client = genai.Client(api_key=api_key)
available_models = list_gemini_models()

# --- Sidebar Inputs (Context) ---
st.sidebar.header("Add Context")
model_choice = st.sidebar.selectbox("Choose a Model", available_models)

uploaded_file = st.sidebar.file_uploader(
    "Upload a file", 
    type=["pdf", "png", "jpg", "jpeg", "mp3", "wav", "mp4"],
    on_change=lambda: st.session_state.pop('context', None) # Clear context on new upload
)
youtube_url = st.sidebar.text_input(
    "YouTube URL", 
    on_change=lambda: st.session_state.pop('context', None)
)
# camera_img = st.sidebar.camera_input(
#     "Take a Photo", 
#     on_change=lambda: st.session_state.pop('context', None)
# )

# Process and store context
if uploaded_file and not st.session_state.context:
    file_type = uploaded_file.type
    st.session_state.context = {"type": file_type, "data": uploaded_file, "name": uploaded_file.name}
    st.sidebar.success(f"Added {uploaded_file.name} to context.")
elif youtube_url and not st.session_state.context:
    st.session_state.context = {"type": "youtube", "data": youtube_url, "name": youtube_url}
    st.sidebar.success("Added YouTube URL to context.")
    st.sidebar.video(youtube_url)
# elif camera_img and not st.session_state.context:
#     st.session_state.context = {"type": "image", "data": camera_img, "name": "camera_photo.jpg"}
#     st.sidebar.success("Added photo to context.")
#     st.sidebar.image(camera_img)

# --- Sidebar Chat Management ---
st.sidebar.header("Chat Management")
if st.sidebar.button("Clear Chat History"):
    st.session_state.messages = [{"role": "assistant", "content": "How can I help you today?"}]
    st.session_state.context = None
    st.rerun()

st.sidebar.subheader("Download Chat")
chat_history_md = format_chat_history(st.session_state.messages)
download_format = st.sidebar.radio("Format:", ("TXT", "Markdown", "DOCX", "PPTX"))

now = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M")
file_name_base = f"gemini_chat_{now}"

if download_format == "TXT":
    st.sidebar.download_button(
        label="Download TXT",
        data=chat_history_md.encode("utf-8"),
        file_name=f"{file_name_base}.txt",
        mime="text/plain",
    )
elif download_format == "Markdown":
    st.sidebar.download_button(
        label="Download MD",
        data=chat_history_md.encode("utf-8"),
        file_name=f"{file_name_base}.md",
        mime="text/markdown",
    )
# elif download_format == "PDF":
#     pdf_bytes = convert_md_to_pdf(chat_history_md)
#     st.sidebar.download_button(
#         label="Download PDF",
#         data=pdf_bytes,
#         file_name=f"{file_name_base}.pdf",
#         mime="application/pdf",
#     )
elif download_format == "DOCX":
    docx_bytes = convert_md_to_docx(chat_history_md)
    st.sidebar.download_button(
        label="Download DOCX",
        data=docx_bytes,
        file_name=f"{file_name_base}.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
elif download_format == "PPTX":
    pptx_bytes = convert_md_to_pptx(chat_history_md)
    st.sidebar.download_button(
        label="Download PPTX",
        data=pptx_bytes,
        file_name=f"{file_name_base}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# --- Main Chat Interface ---
st.title("💬 Chat with Google Gemini")

# Display chat messages
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# Chat input
if prompt := st.chat_input("Ask about text, your upload, the YouTube video, or your photo..."):
    
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    contents = [prompt]
    context_used = False

    # Check for and prepare context
    if st.session_state.context:
        context = st.session_state.context
        context_type = context["type"]
        context_data = context["data"]
        
        try:
            if "pdf" in context_type:
                pdf_text = process_pdf(context_data.getvalue())
                if pdf_text:
                    contents.append(pdf_text)
                    context_used = True
            
            elif "image" in context_type:
                contents.append(types.Part.from_bytes(data=context_data.getvalue(), mime_type=context_type))
                context_used = True

            elif "audio" in context_type or "video" in context_type:
                contents.append(types.Part.from_bytes(data=context_data.getvalue(), mime_type=context_type))
                context_used = True
            
            elif context_type == "youtube":
                with st.spinner("Processing YouTube video..."):
                    audio_stream = get_youtube_stream(context_data, "audio")
                    if audio_stream:
                        buffer = io.BytesIO()
                        audio_stream.stream_to_buffer(buffer)
                        buffer.seek(0)
                        contents.append(types.Part.from_bytes(data=buffer.read(), mime_type="audio/mpeg"))
                        context_used = True
            
            if context_used:
                with st.chat_message("assistant"):
                    st.info(f"Using context: **{context['name']}**")
        
        except Exception as e:
            st.error(f"Error processing context: {e}")

    # Generate and display response
    with st.chat_message("assistant"):
        with st.spinner("Gemini is thinking..."):
            response = generate_gemini_content(client, model_choice, contents)
            if response:
                st.markdown(response.text)
                # Add only the AI's text response to session state
                st.session_state.messages.append({"role": "assistant", "content": response.text})
            else:
                st.error("Failed to get a response from Gemini.")

    # Clear context after use
    if context_used:
        st.session_state.context = None